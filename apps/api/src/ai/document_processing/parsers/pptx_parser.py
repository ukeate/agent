"""PowerPoint文档解析器"""

from pathlib import Path
from typing import List, Dict, Any
import base64
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
from .base_parser import BaseParser, ParsedDocument, ParsedElement

logger = get_logger(__name__)

class PowerPointParser(BaseParser):
    """PowerPoint文档解析器
    
    支持.pptx格式的幻灯片文本、图像、表格提取
    """
    
    SUPPORTED_EXTENSIONS = [".pptx", ".ppt"]
    
    def __init__(
        self, 
        extract_images: bool = True,
        extract_notes: bool = True,
        extract_layouts: bool = True
    ):
        """初始化PowerPoint解析器
        
        Args:
            extract_images: 是否提取图像
            extract_notes: 是否提取演讲者备注
            extract_layouts: 是否提取布局信息
        """
        super().__init__()
        self.extract_images = extract_images
        self.extract_notes = extract_notes
        self.extract_layouts = extract_layouts
    
    async def parse(self, file_path: Path) -> ParsedDocument:
        """解析PowerPoint文档
        
        Args:
            file_path: PowerPoint文件路径
            
        Returns:
            解析后的文档
        """
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # 检查文件格式
        if file_path.suffix.lower() == ".ppt":
            logger.warning("Old .ppt format detected. Consider converting to .pptx for better support")
        
        doc_id = self.generate_doc_id(file_path)
        metadata = self.extract_metadata(file_path)
        elements = []
        
        try:
            # 打开演示文稿
            prs = Presentation(str(file_path))
            
            # 提取演示文稿属性
            core_props = prs.core_properties
            metadata.update({
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "keywords": core_props.keywords or "",
                "created": core_props.created.isoformat() if core_props.created else "",
                "modified": core_props.modified.isoformat() if core_props.modified else "",
                "slide_count": len(prs.slides),
            })
            
            # 解析每个幻灯片
            for slide_idx, slide in enumerate(prs.slides):
                slide_elements = await self._parse_slide(slide, slide_idx + 1)
                elements.extend(slide_elements)
                
                # 提取演讲者备注
                if self.extract_notes:
                    notes = await self._extract_notes(slide, slide_idx + 1)
                    if notes:
                        elements.append(notes)
            
        except Exception as e:
            logger.error(f"Error parsing PowerPoint document: {e}")
            raise
        
        return ParsedDocument(
            doc_id=doc_id,
            file_path=str(file_path),
            file_type="powerpoint",
            elements=elements,
            metadata=metadata
        )
    
    async def _parse_slide(
        self, 
        slide, 
        slide_number: int
    ) -> List[ParsedElement]:
        """解析单个幻灯片
        
        Args:
            slide: 幻灯片对象
            slide_number: 幻灯片编号
            
        Returns:
            解析的元素列表
        """
        elements = []
        
        # 提取幻灯片标题
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                elements.append(ParsedElement(
                    content=title_text,
                    element_type="heading",
                    metadata={
                        "slide_number": slide_number,
                        "shape_type": "title",
                        "heading_level": 1,
                    }
                ))
        
        # 遍历所有形状
        for shape in slide.shapes:
            # 跳过标题（已处理）
            if shape == slide.shapes.title:
                continue
            
            # 文本框
            if shape.has_text_frame:
                text_element = await self._parse_text_frame(shape, slide_number)
                if text_element:
                    elements.append(text_element)
            
            # 表格
            elif shape.has_table:
                table_element = await self._parse_table(shape.table, slide_number)
                elements.append(table_element)
            
            # 图片
            elif self.extract_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_element = await self._parse_image(shape, slide_number)
                if image_element:
                    elements.append(image_element)
            
            # 图表
            elif shape.has_chart:
                chart_element = await self._parse_chart(shape, slide_number)
                elements.append(chart_element)
        
        # 添加布局信息
        if self.extract_layouts and elements:
            elements[0].metadata["layout_name"] = slide.slide_layout.name
        
        return elements
    
    async def _parse_text_frame(self, shape, slide_number: int) -> ParsedElement:
        """解析文本框
        
        Args:
            shape: 形状对象
            slide_number: 幻灯片编号
            
        Returns:
            解析的元素或None
        """
        text_parts = []
        bullet_levels = []
        
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                text_parts.append(text)
                # 记录项目符号层级
                if paragraph.level is not None:
                    bullet_levels.append(paragraph.level)
        
        if not text_parts:
            return None
        
        content = "\n".join(text_parts)
        
        # 判断是否为列表
        element_type = "list" if bullet_levels else "text"
        
        metadata = {
            "slide_number": slide_number,
            "shape_type": "text_frame",
        }
        
        if bullet_levels:
            metadata["bullet_levels"] = bullet_levels
        
        return ParsedElement(
            content=content,
            element_type=element_type,
            metadata=metadata
        )
    
    async def _parse_table(self, table, slide_number: int) -> ParsedElement:
        """解析表格
        
        Args:
            table: 表格对象
            slide_number: 幻灯片编号
            
        Returns:
            解析的元素
        """
        # 提取表格数据
        table_data = []
        headers = []
        
        for row_idx, row in enumerate(table.rows):
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                row_data.append(cell_text)
            
            if row_idx == 0:
                headers = row_data
            else:
                table_data.append(row_data)
        
        # 转换为文本格式
        content_lines = []
        
        # 添加表头
        if headers:
            content_lines.append(" | ".join(headers))
            content_lines.append("-" * (len(" | ".join(headers))))
        
        # 添加数据行
        for row in table_data:
            content_lines.append(" | ".join(row))
        
        content = "\n".join(content_lines)
        
        return ParsedElement(
            content=content,
            element_type="table",
            metadata={
                "slide_number": slide_number,
                "shape_type": "table",
                "headers": headers,
                "rows": table_data,
                "num_rows": len(table.rows),
                "num_columns": len(table.columns),
            }
        )
    
    async def _parse_image(self, shape, slide_number: int) -> ParsedElement:
        """解析图像
        
        Args:
            shape: 形状对象
            slide_number: 幻灯片编号
            
        Returns:
            解析的元素或None
        """
        try:
            # 获取图像数据
            image = shape.image
            image_bytes = image.blob
            
            # 获取图像信息
            img = Image.open(io.BytesIO(image_bytes))
            
            # 转换为base64
            img_base64 = base64.b64encode(image_bytes).decode()
            
            # 尝试获取图像描述（替代文本）
            alt_text = ""
            if hasattr(shape, 'name'):
                alt_text = shape.name
            
            return ParsedElement(
                content=img_base64,
                element_type="image",
                metadata={
                    "slide_number": slide_number,
                    "shape_type": "picture",
                    "width": img.width,
                    "height": img.height,
                    "format": img.format,
                    "alt_text": alt_text,
                }
            )
            
        except Exception as e:
            logger.warning(f"Error extracting image: {e}")
            return None
    
    async def _parse_chart(self, shape, slide_number: int) -> ParsedElement:
        """解析图表
        
        Args:
            shape: 形状对象
            slide_number: 幻灯片编号
            
        Returns:
            解析的元素
        """
        chart_info = {
            "title": shape.chart.chart_title.text_frame.text if shape.chart.has_title else "Untitled Chart",
            "chart_type": str(shape.chart.chart_type),
        }
        
        # 尝试提取图表数据
        try:
            # 提取系列数据
            series_info = []
            for series in shape.chart.series:
                series_info.append({
                    "name": series.name,
                    "values_count": len(series.values) if hasattr(series, 'values') else 0
                })
            
            chart_info["series"] = series_info
            chart_info["series_count"] = len(series_info)
            
        except Exception as e:
            logger.warning(f"Error extracting chart data: {e}")
        
        content = f"Chart: {chart_info['title']} (Type: {chart_info['chart_type']})"
        
        return ParsedElement(
            content=content,
            element_type="chart",
            metadata={
                "slide_number": slide_number,
                "shape_type": "chart",
                **chart_info
            }
        )
    
    async def _extract_notes(self, slide, slide_number: int) -> ParsedElement:
        """提取演讲者备注
        
        Args:
            slide: 幻灯片对象
            slide_number: 幻灯片编号
            
        Returns:
            备注元素或None
        """
        if not slide.has_notes_slide:
            return None
        
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text.strip()
        
        if not notes_text:
            return None
        
        return ParsedElement(
            content=notes_text,
            element_type="notes",
            metadata={
                "slide_number": slide_number,
                "shape_type": "speaker_notes",
            }
        )
from src.core.logging import get_logger
